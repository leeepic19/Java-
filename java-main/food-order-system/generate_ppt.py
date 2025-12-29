from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import os

def create_ppt(markdown_file, output_file):
    prs = Presentation()

    if not os.path.exists(markdown_file):
        print(f"Error: File not found: {markdown_file}")
        return

    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()

    lines = content.split('\n')
    
    current_slide = None
    body_shape = None
    tf = None
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('# '):
            # Main presentation title, ignore as we use Slide 1
            pass
            
        elif line.startswith('## Slide'):
            # New Slide
            if 'Slide 1:' in line or '封面' in line:
                slide_layout = prs.slide_layouts[0] # Title Slide
            else:
                slide_layout = prs.slide_layouts[1] # Title and Content

            current_slide = prs.slides.add_slide(slide_layout)
            
            # Set Title
            if ':' in line:
                title = line.split(':', 1)[1].strip()
            else:
                title = line.replace('##', '').strip()
                
            title_shape = current_slide.shapes.title
            title_shape.text = title
            
            # Prepare body
            if slide_layout == prs.slide_layouts[1]:
                body_shape = current_slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear() # Clear default text
            else:
                # Title slide usually has title and subtitle
                # We will put content in subtitle placeholder if available
                if len(current_slide.placeholders) > 1:
                    body_shape = current_slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()
                else:
                    tf = None

        elif line.startswith('- '):
            if current_slide and tf:
                p = tf.add_paragraph()
                text = line[2:]
                
                # Simple bold handling
                parts = re.split(r'(\*\*.*?\*\*)', text)
                for part in parts:
                    run = p.add_run()
                    if part.startswith('**') and part.endswith('**'):
                        run.text = part[2:-2]
                        run.font.bold = True
                    else:
                        run.text = part
                
                p.level = 0

        elif line.startswith('    - '):
            if current_slide and tf:
                p = tf.add_paragraph()
                text = line[6:]
                parts = re.split(r'(\*\*.*?\*\*)', text)
                for part in parts:
                    run = p.add_run()
                    if part.startswith('**') and part.endswith('**'):
                        run.text = part[2:-2]
                        run.font.bold = True
                    else:
                        run.text = part
                p.level = 1

    prs.save(output_file)
    print(f"PPT generated: {output_file}")

if __name__ == "__main__":
    md_path = r'c:\Users\35278\Desktop\java-main\java-main\food-order-system\项目技术报告.md'
    ppt_path = r'c:\Users\35278\Desktop\java-main\java-main\food-order-system\项目汇报.pptx'
    create_ppt(md_path, ppt_path)
